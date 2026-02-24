#!/usr/bin/env python3
"""Meeting2Deck: Generate slides.pptx from meeting analysis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "output", "slides.pptx")

# Ensure output directory exists
os.makedirs(os.path.join(OUTPUT_DIR, "output"), exist_ok=True)

# Color palette
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MEDIUM = RGBColor(0x66, 0x66, 0x66)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
BOX_BG = RGBColor(0xE8, 0xF0, 0xFE)
ARROW_COLOR = RGBColor(0x5B, 0x9B, 0xD5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=16,
                     font_color=TEXT_DARK, line_spacing=1.5):
    """Add bullet points to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Arial"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_title_bar(slide, title_text):
    """Add a consistent title bar at the top of content slides."""
    # Title background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()

    # Title text
    add_textbox(slide, 0.8, 0.25, 11, 0.7, title_text,
                font_size=28, font_color=WHITE, bold=True)


def add_rounded_box(slide, left, top, width, height, text,
                    fill_color=BOX_BG, text_color=TEXT_DARK, font_size=14,
                    bold=False):
    """Add a rounded rectangle box with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Arial"

    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a connector arrow."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = ARROW_COLOR
    connector.line.width = Pt(2)
    return connector


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, DARK_NAVY)

# Accent line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1.5), Inches(2.8),
    Inches(2), Inches(0.06)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_textbox(slide, 1.5, 1.5, 10, 1.2,
            "컨택센터 시스템 아키텍처 구성 회의",
            font_size=36, font_color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.0, 10, 0.6,
            "IPCC + AICC 통합 아키텍처 검토",
            font_size=20, font_color=RGBColor(0xA0, 0xC4, 0xE8))
add_textbox(slide, 1.5, 5.5, 10, 0.5,
            "2026년 2월 24일",
            font_size=16, font_color=RGBColor(0x99, 0x99, 0x99))

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Executive Summary")

bullets = [
    "LG U+ 기반 IPCC + AICC 통합 컨택센터 아키텍처 구성 검토",
    "핵심 인프라: PBX → IVR → CTI 흐름 + Cloud 연동",
    "AI 고객 응대(AICC) 기능 추가를 통한 CS 고도화 방향 확인",
    "S1 시스템 상세 구성 확정 필요 (후속 회의 예정)",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# ============================================================
# SLIDE 3: Background / Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Background / Problem Definition")

bullets = [
    "기존 IPCC(IP Contact Center) 시스템의 한계 대응 필요 [추정]",
    "고객 응대 자동화 및 AI 기반 서비스 고도화 요구 증가",
    "Cloud 연동을 통한 확장성 및 유연성 확보 필요",
    "LG U+ 솔루션 기반 통합 아키텍처 설계 검토 착수",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# ============================================================
# SLIDE 4: Current State Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Current State Analysis")

# Left column
add_textbox(slide, 1.0, 1.8, 5.5, 0.5, "현재 인프라 구성요소",
            font_size=20, font_color=MEDIUM_BLUE, bold=True)
left_bullets = [
    "PBX (Private Branch Exchange) — 전화 교환 시스템",
    "IVR (Interactive Voice Response) — 자동 음성 응답",
    "CTI (Computer Telephony Integration) — 전화-컴퓨터 통합",
]
add_bullet_slide(slide, 1.0, 2.5, 5.5, 3.0, left_bullets, font_size=16)

# Right column
add_textbox(slide, 7.0, 1.8, 5.5, 0.5, "전환 방향",
            font_size=20, font_color=MEDIUM_BLUE, bold=True)
right_bullets = [
    "IPCC + AICC 통합 솔루션 도입",
    "LG U+ 솔루션 파트너십 활용",
    "Cloud 연동으로 인프라 확장",
]
add_bullet_slide(slide, 7.0, 2.5, 5.5, 3.0, right_bullets, font_size=16)

# ============================================================
# SLIDE 5: Key Discussion Points
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Key Discussion Points")

bullets = [
    "S1 시스템 컨택센터 아키텍처 구성 방향 논의",
    "PBX-IVR-CTI 기반 기존 인프라 활용 방안",
    "AICC(AI Contact Center) 기능 통합 범위 검토",
    "Cloud 연동 시 인프라 구성 변경 사항 확인",
    "LG U+ 솔루션 기술 적합성 초기 평가",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# ============================================================
# SLIDE 6: Strategic Direction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Strategic Direction")

bullets = [
    "IPCC → IPCC + AICC 통합 전환으로 CS 경쟁력 강화 [추정]",
    "기존 PBX-IVR-CTI 인프라 유지하며 AI 레이어 추가",
    "Cloud 기반 유연한 확장 아키텍처 구현",
    "LG U+ 솔루션 활용으로 구축 리스크 최소화 [추정]",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# ============================================================
# SLIDE 7: Architecture Diagram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "System Architecture")

# Top-level: LG U+ (IPCC + AICC) banner
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.5), Inches(1.5),
    Inches(6.3), Inches(0.8)
)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_NAVY
shape.line.fill.background()
tf = shape.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "LG U+  |  IPCC + AICC"
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Arial"

# PBX Box
add_rounded_box(slide, 2.0, 3.3, 2.2, 1.0, "PBX\nPrivate Branch Exchange",
                fill_color=RGBColor(0xDB, 0xEA, 0xFE), font_size=13, bold=True)

# IVR Box
add_rounded_box(slide, 5.5, 3.3, 2.2, 1.0, "IVR\nInteractive Voice Response",
                fill_color=RGBColor(0xDB, 0xEA, 0xFE), font_size=13, bold=True)

# CTI Box
add_rounded_box(slide, 3.8, 5.2, 2.2, 1.0, "CTI\nComputer Telephony Integration",
                fill_color=RGBColor(0xDB, 0xEA, 0xFE), font_size=13, bold=True)

# Cloud Box
add_rounded_box(slide, 9.0, 3.3, 2.2, 1.0, "Cloud\nInfrastructure",
                fill_color=RGBColor(0xE8, 0xF5, 0xE9), text_color=RGBColor(0x2E, 0x7D, 0x32),
                font_size=13, bold=True)

# Arrows: PBX → IVR
add_arrow(slide, 4.2, 3.8, 5.5, 3.8)

# Arrows: PBX → CTI
add_arrow(slide, 3.1, 4.3, 4.3, 5.2)

# Arrows: IVR → CTI
add_arrow(slide, 6.2, 4.3, 5.5, 5.2)

# Arrows: IVR → Cloud
add_arrow(slide, 7.7, 3.8, 9.0, 3.8)

# Label
add_textbox(slide, 1.0, 6.5, 11, 0.5,
            "* PBX가 수신 호를 IVR로 라우팅하고, CTI가 통화-데이터 통합 처리. Cloud로 확장 연동.",
            font_size=12, font_color=TEXT_MEDIUM)

# ============================================================
# SLIDE 8: Risks & Considerations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Risks & Considerations")

bullets = [
    "온프레미스-Cloud 하이브리드 구성 시 보안 및 안정성 리스크 [추정]",
    "AICC 도입 시 기존 IPCC 운영 프로세스 영향 검토 필요",
    "LG U+ 솔루션 종속성(Vendor Lock-in) 리스크 평가 필요 [추정]",
    "시스템 전환 기간 중 서비스 연속성 확보 방안 수립 필요",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# ============================================================
# SLIDE 9: Action Items
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Action Items")

items = [
    "① S1 시스템 구성도 상세화 — 담당자 지정 필요",
    "② IPCC-AICC 통합 방안 구체화 — 기술 검토 포함",
    "③ Cloud 연동 아키텍처 설계 — 보안 요건 포함",
    "④ LG U+ 솔루션 기술 스펙 검토 및 PoC 계획 수립",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, items, font_size=20)

# ============================================================
# SLIDE 10: Roadmap / Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title_bar(slide, "Roadmap / Next Steps")

bullets = [
    "단기: S1 시스템 상세 구성안 확정 및 기술 검토",
    "중기: IPCC + AICC 통합 PoC 수행 및 결과 평가 [추정]",
    "장기: Cloud 연동 포함 전체 아키텍처 구축 완료 [추정]",
]
add_bullet_slide(slide, 1.0, 1.8, 11, 4.5, bullets, font_size=20)

# Next meeting prep
add_textbox(slide, 1.0, 4.5, 11, 0.5, "다음 회의 준비사항",
            font_size=20, font_color=MEDIUM_BLUE, bold=True)
prep_bullets = [
    "S1 시스템 상세 구성안 사전 작성",
    "LG U+ IPCC/AICC 기술 스펙 자료 수집",
    "Cloud 연동 보안/안정성 검토 사항 정리",
]
add_bullet_slide(slide, 1.0, 5.1, 11, 2.0, prep_bullets, font_size=16,
                 font_color=TEXT_MEDIUM)

# Save
prs.save(OUTPUT_PATH)
print(f"Slides saved to: {OUTPUT_PATH}")
