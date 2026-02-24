"""Meeting2Deck 프레젠테이션 디자인 템플릿.

Claude CLI 에이전트가 이 모듈을 import하여 컨설팅 수준의 PPTX를 생성합니다.

사용법:
    import sys, os
    sys.path.insert(0, os.getcwd())
    from slide_template import DeckBuilder

    deck = DeckBuilder("IPCC + AICC 아키텍처", date="2026.02.24", org="팀명")

    deck.add_title_slide(
        subtitle="통합 솔루션 아키텍처",
        description="PBX · IVR · CTI 기반 차세대 통합 아키텍처"
    )
    deck.add_content_slide("01", "Executive Summary", "핵심 요약", [
        "첫번째 불릿 포인트",
        "두번째 불릿 포인트",
    ])
    deck.add_cards_slide("02", "Background", "배경 및 문제", [
        {"icon": "1", "title": "과제", "body": "설명..."},
        {"icon": "2", "title": "방향", "body": "설명..."},
        {"icon": "3", "title": "목표", "body": "설명..."},
    ])
    deck.add_diagram_slide("07", "Architecture", "시스템 구성도", [
        {"name": "PBX", "desc": "전화 교환기"},
        {"name": "IVR", "desc": "자동 응답"},
        {"name": "CTI", "desc": "컴퓨터-전화 연동"},
        {"name": "Cloud", "desc": "AI 서비스"},
    ])
    deck.add_table_slide("09", "Action Items", "액션 아이템",
        ["#", "담당", "내용", "기한", "상태"],
        [["1", "TBD", "구성도 상세 설계", "TBD", "미착수"]]
    )
    deck.add_closing_slide(
        message="Thank You",
        submessage="차세대 컨택센터의 새로운 기준을 만들겠습니다."
    )
    deck.save("output/slides.pptx")
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os


def _rgb(h):
    """'#RRGGBB' → RGBColor"""
    h = h.lstrip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class DeckBuilder:
    """컨설팅 수준의 프레젠테이션 빌더.

    16:9 와이드스크린, 카드 기반 레이아웃, 체계적 색상 팔레트를 적용합니다.
    """

    # ── 슬라이드 크기 (16:9) ──
    SW, SH = 9144000, 5143500  # 10" x 5.63"

    # ── 레이아웃 상수 ──
    MX = 640080       # 좌우 마진 (0.7")
    CW = 7863840      # 콘텐츠 폭 (SW - 2*MX = 8.6")
    Y_SEC = 365760    # 섹션 헤더
    Y_TTL = 822960    # 페이지 제목
    Y_DESC = 1371600  # 설명
    Y_BODY = 1920240  # 본문 시작
    Y_FT = 4800600    # 푸터
    H_FT = 342900     # 푸터 높이

    # ── 컬러 팔레트 ──
    C = {
        'accent':   _rgb('#0070C0'),   # 주 강조 (파랑)
        'accent2':  _rgb('#00B4D8'),   # 보조 강조 (시안)
        'highlight': _rgb('#E63946'),  # 하이라이트 (레드)
        'bg_dark':  _rgb('#0A1628'),   # 타이틀/클로징 배경
        'bg_light': _rgb('#F0F4F8'),   # 콘텐츠 배경
        'card':     _rgb('#FFFFFF'),   # 카드 배경
        'text':     _rgb('#1E293B'),   # 본문 (진한)
        'text2':    _rgb('#475569'),   # 본문 (중간)
        'text3':    _rgb('#94A3B8'),   # 본문 (연한)
        'white':    _rgb('#FFFFFF'),
        'footer':   _rgb('#0F1B2D'),   # 푸터 배경
        'divider':  _rgb('#E2E8F0'),   # 구분선
    }
    ICON_C = [_rgb(c) for c in [
        '#0070C0', '#00B4D8', '#38A169', '#DD6B20', '#805AD5', '#E53E3E',
    ]]

    FONT = 'Apple SD Gothic Neo'

    # ─────────────────────────────────────────────
    def __init__(self, title, date='', org=''):
        self.prs = Presentation()
        self.prs.slide_width = Emu(self.SW)
        self.prs.slide_height = Emu(self.SH)
        self.title = title
        self.date = date
        self.org = org
        self._pg = 0

    # ── 저수준 헬퍼 ──

    def _slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

    def _rect(self, s, x, y, w, h, c=None):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if c:
            sh.fill.solid()
            sh.fill.fore_color.rgb = c
        else:
            sh.fill.background()
        sh.line.fill.background()
        return sh

    def _rrect(self, s, x, y, w, h, c=None):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        if c:
            sh.fill.solid()
            sh.fill.fore_color.rgb = c
        else:
            sh.fill.background()
        sh.line.fill.background()
        return sh

    def _oval(self, s, x, y, d, c):
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()
        return sh

    def _txt(self, s, x, y, w, h, text, sz=11, c=None, b=False,
             align=PP_ALIGN.LEFT, va=None):
        """텍스트 박스 추가."""
        sh = s.shapes.add_textbox(x, y, w, h)
        tf = sh.text_frame
        tf.word_wrap = True
        if va:
            tf.vertical_anchor = va
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.name = self.FONT
        if c:
            run.font.color.rgb = c
        return sh

    # ── 공통 슬라이드 요소 ──

    def _bar(self, s):
        """상단 악센트 바 (3pt 파랑 라인)."""
        self._rect(s, 0, 0, self.SW, Pt(3), self.C['accent'])

    def _footer(self, s):
        """하단 푸터 바 + 페이지 번호."""
        self._pg += 1
        self._rect(s, 0, self.Y_FT, self.SW, self.H_FT, self.C['footer'])
        self._txt(s, self.MX, self.Y_FT, Emu(4572000), self.H_FT,
                  self.title, sz=8, c=self.C['text3'])
        self._txt(s, Emu(self.SW - self.MX - 640080), self.Y_FT,
                  Emu(640080), self.H_FT,
                  str(self._pg), sz=8, c=self.C['text3'], align=PP_ALIGN.RIGHT)

    def _sec(self, s, num, name):
        """섹션 인디케이터 (컬러 바 + 번호/이름)."""
        self._rect(s, self.MX, self.Y_SEC, Emu(54864), Emu(320040), self.C['accent'])
        self._txt(s, self.MX + Emu(182880), self.Y_SEC, Emu(4572000), Emu(320040),
                  f"{num}  {name}", sz=11, c=self.C['accent'], b=True)

    def _ttl(self, s, text):
        """페이지 제목 (26pt Bold)."""
        self._txt(s, self.MX, self.Y_TTL, Emu(self.CW), Emu(502920),
                  text, sz=26, c=self.C['text'], b=True)

    def _desc(self, s, text):
        """페이지 설명 (12pt)."""
        self._txt(s, self.MX, self.Y_DESC, Emu(self.CW), Emu(365760),
                  text, sz=12, c=self.C['text2'])

    def _bg(self, s, c):
        """슬라이드 배경색 설정."""
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = c

    # ════════════════════════════════════════════════
    #  슬라이드 타입
    # ════════════════════════════════════════════════

    def add_title_slide(self, subtitle='', description=''):
        """표지 슬라이드 — 어두운 배경, 대형 타이틀."""
        s = self._slide()
        self._bg(s, self.C['bg_dark'])
        self._bar(s)

        # 악센트 라인
        self._rect(s, self.MX, Emu(1600200), Emu(1828800), Pt(4), self.C['accent'])

        # 메인 타이틀
        self._txt(s, self.MX, Emu(1783080), Emu(self.CW), Emu(960120),
                  self.title, sz=44, c=self.C['white'], b=True)

        if subtitle:
            self._txt(s, self.MX, Emu(2697480), Emu(self.CW), Emu(594360),
                      subtitle, sz=24, c=self.C['accent2'], b=True)

        if description:
            self._txt(s, self.MX, Emu(3429000), Emu(5943600), Emu(731520),
                      description, sz=13, c=self.C['text3'])

        info = ' | '.join(filter(None, [self.org, self.date]))
        if info:
            self._txt(s, self.MX, Emu(4114800), Emu(self.CW), Emu(365760),
                      info, sz=11, c=self.C['text3'])

        # 표지 푸터 (번호 없음)
        self._rect(s, 0, self.Y_FT, self.SW, self.H_FT, self.C['footer'])
        ft = ' | '.join(filter(None, ['Confidential', self.org, self.date]))
        self._txt(s, self.MX, self.Y_FT, Emu(7680960), self.H_FT,
                  ft, sz=8.5, c=self.C['text3'])

    def add_content_slide(self, num, section, title, bullets, description=''):
        """불릿 포인트 슬라이드 — 흰 카드 위에 불릿 리스트.

        Args:
            num: 섹션 번호 (예: "01")
            section: 섹션 이름 (예: "Executive Summary")
            title: 슬라이드 제목
            bullets: 불릿 텍스트 리스트
            description: (선택) 제목 아래 설명 문구
        """
        s = self._slide()
        self._bg(s, self.C['bg_light'])
        self._bar(s)
        self._sec(s, num, section)
        self._ttl(s, title)
        if description:
            self._desc(s, description)

        # 카드 영역 계산
        cy = self.Y_DESC + Emu(457200) if description else self.Y_TTL + Emu(594360)
        ch = self.Y_FT - cy - Emu(137160)
        card = self._rrect(s, self.MX, cy, Emu(self.CW), ch, self.C['card'])

        # 불릿 텍스트
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(274320)
        tf.margin_right = Emu(274320)
        tf.margin_top = Emu(228600)

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(12)
            run.font.name = self.FONT
            run.font.color.rgb = self.C['text2']

        self._footer(s)

    def add_cards_slide(self, num, section, title, cards, description=''):
        """카드 그리드 슬라이드 — 2~4장 카드를 나란히 배치.

        Args:
            cards: [{"icon": "1", "title": "제목", "body": "설명"}, ...]
                icon은 카드 원 안에 표시할 텍스트 (숫자, 이모지, 약어 등)
        """
        s = self._slide()
        self._bg(s, self.C['bg_light'])
        self._bar(s)
        self._sec(s, num, section)
        self._ttl(s, title)
        if description:
            self._desc(s, description)

        n = min(len(cards), 4)
        if n == 0:
            self._footer(s)
            return

        gap = 228600
        cy = self.Y_DESC + Emu(457200) if description else self.Y_TTL + Emu(594360)
        ch = self.Y_FT - cy - Emu(137160)
        cw = (self.CW - gap * (n - 1)) // n

        for i, c in enumerate(cards[:4]):
            x = self.MX + i * (cw + gap)
            ic = self.ICON_C[i % len(self.ICON_C)]

            # 카드 배경
            self._rrect(s, x, cy, cw, ch, self.C['card'])

            # 아이콘 원
            d = Emu(411480)
            ox, oy = x + Emu(182880), cy + Emu(182880)
            circle = self._oval(s, ox, oy, d, ic)

            # 아이콘 텍스트
            ctf = circle.text_frame
            ctf.word_wrap = False
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = c.get('icon', str(i + 1))
            cr.font.size = Pt(16)
            cr.font.color.rgb = self.C['white']
            cr.font.bold = True
            cr.font.name = self.FONT

            # 카드 제목
            self._txt(s, ox, oy + d + Emu(137160), cw - Emu(365760), Emu(320040),
                      c.get('title', ''), sz=14, c=self.C['text'], b=True)

            # 카드 본문
            body_y = oy + d + Emu(502920)
            body_h = ch - (body_y - cy) - Emu(137160)
            self._txt(s, ox, body_y, cw - Emu(365760), body_h,
                      c.get('body', ''), sz=10, c=self.C['text2'])

        self._footer(s)

    def add_two_column_slide(self, num, section, title, left_bullets, right_bullets,
                             left_title='', right_title='', description=''):
        """2열 슬라이드 — 좌우 카드에 각각 불릿 리스트.

        Args:
            left_bullets: 왼쪽 카드 불릿 리스트
            right_bullets: 오른쪽 카드 불릿 리스트
            left_title: 왼쪽 카드 소제목
            right_title: 오른쪽 카드 소제목
        """
        s = self._slide()
        self._bg(s, self.C['bg_light'])
        self._bar(s)
        self._sec(s, num, section)
        self._ttl(s, title)
        if description:
            self._desc(s, description)

        gap = 228600
        cy = self.Y_DESC + Emu(457200) if description else self.Y_TTL + Emu(594360)
        ch = self.Y_FT - cy - Emu(137160)
        cw = (self.CW - gap) // 2

        for col, (col_title, bullets) in enumerate([
            (left_title, left_bullets), (right_title, right_bullets)
        ]):
            x = self.MX + col * (cw + gap)
            card = self._rrect(s, x, cy, cw, ch, self.C['card'])

            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(182880)
            tf.margin_right = Emu(182880)
            tf.margin_top = Emu(137160)

            # 카드 소제목
            if col_title:
                p = tf.paragraphs[0]
                p.space_after = Pt(12)
                r = p.add_run()
                r.text = col_title
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.name = self.FONT
                r.font.color.rgb = self.C['text']

            for i, bullet in enumerate(bullets):
                if i == 0 and not col_title:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(8)
                r = p.add_run()
                r.text = f"•  {bullet}"
                r.font.size = Pt(11)
                r.font.name = self.FONT
                r.font.color.rgb = self.C['text2']

        self._footer(s)

    def add_table_slide(self, num, section, title, headers, rows):
        """표 슬라이드 — 헤더 행(파랑) + 데이터 행(흰색).

        Args:
            headers: 열 헤더 리스트
            rows: 행 데이터 리스트 (각 행은 리스트)
        """
        s = self._slide()
        self._bg(s, self.C['bg_light'])
        self._bar(s)
        self._sec(s, num, section)
        self._ttl(s, title)

        ty = self.Y_TTL + Emu(594360)
        th = self.Y_FT - ty - Emu(137160)
        nr, nc = len(rows) + 1, len(headers)

        tbl_sh = s.shapes.add_table(nr, nc, self.MX, ty, Emu(self.CW), th)
        tbl = tbl_sh.table

        # 헤더 행
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.C['accent']
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.bold = True
                    r.font.color.rgb = self.C['white']
                    r.font.name = self.FONT

        # 데이터 행 (줄무늬)
        for i, row in enumerate(rows):
            bg = self.C['card'] if i % 2 == 0 else self.C['bg_light']
            for j, v in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.text = str(v)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)
                        r.font.name = self.FONT
                        r.font.color.rgb = self.C['text2']

        self._footer(s)

    def add_diagram_slide(self, num, section, title, nodes):
        """아키텍처/플로우 다이어그램 슬라이드 — 노드를 수평 배치하고 화살표 연결.

        Args:
            nodes: [{"name": "PBX", "desc": "전화 교환기"}, ...]
                선택적 키: "icon" (원 안 텍스트, 기본값은 name 첫 2글자)
        """
        s = self._slide()
        self._bg(s, self.C['bg_light'])
        self._bar(s)
        self._sec(s, num, section)
        self._ttl(s, title)

        n = len(nodes)
        if n == 0:
            self._footer(s)
            return

        ay = self.Y_TTL + Emu(594360)
        ah = self.Y_FT - ay - Emu(137160)

        # 노드 크기 계산 (노드 수에 따라 조절)
        if n <= 4:
            bw = Emu(1371600)
        elif n <= 6:
            bw = Emu(1143000)
        else:
            bw = Emu(914400)
        bh = Emu(1600200)
        gap = Emu(365760) if n <= 5 else Emu(228600)

        tw = n * bw + (n - 1) * gap
        sx = (self.SW - tw) // 2
        by = ay + (ah - bh) // 2

        for i, nd in enumerate(nodes):
            x = sx + i * (bw + gap)
            ic = self.ICON_C[i % len(self.ICON_C)]

            # 카드
            self._rrect(s, x, by, bw, bh, self.C['card'])

            # 상단 컬러 바
            self._rect(s, x + Emu(4572), by, bw - Emu(9144), Pt(4), ic)

            # 아이콘 원
            d = Emu(365760)
            cx = x + (bw - d) // 2
            circle = self._oval(s, cx, by + Emu(228600), d, ic)

            ctf = circle.text_frame
            ctf.word_wrap = False
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = nd.get('icon', nd.get('name', '')[:2])
            cr.font.size = Pt(13)
            cr.font.color.rgb = self.C['white']
            cr.font.bold = True
            cr.font.name = self.FONT

            # 이름
            self._txt(s, x + Emu(45720), by + Emu(685800),
                      bw - Emu(91440), Emu(274320),
                      nd.get('name', ''), sz=13, c=self.C['text'], b=True,
                      align=PP_ALIGN.CENTER)

            # 설명
            self._txt(s, x + Emu(91440), by + Emu(1005840),
                      bw - Emu(182880), bh - Emu(1143000),
                      nd.get('desc', ''), sz=9, c=self.C['text2'],
                      align=PP_ALIGN.CENTER)

            # 화살표
            if i < n - 1:
                self._txt(s, x + bw, by + bh // 2 - Emu(137160),
                          gap, Emu(274320), '→', sz=28, c=self.C['accent'],
                          align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

        self._footer(s)

    def add_closing_slide(self, message='Thank You', submessage=''):
        """마무리 슬라이드 — 어두운 배경, 감사 메시지."""
        s = self._slide()
        self._bg(s, self.C['bg_dark'])
        self._bar(s)

        # 메인 메시지
        self._txt(s, self.MX, Emu(1371600), Emu(self.CW), Emu(914400),
                  message, sz=50, c=self.C['white'], b=True, align=PP_ALIGN.CENTER)

        # 구분선
        self._rect(s, Emu(3200400), Emu(2377440), Emu(2743200), Pt(3), self.C['accent'])

        if submessage:
            self._txt(s, Emu(1371600), Emu(2651760), Emu(6400800), Emu(822960),
                      submessage, sz=16, c=self.C['accent2'], align=PP_ALIGN.CENTER)

        if self.org:
            self._txt(s, Emu(1828800), Emu(3657600), Emu(5486400), Emu(365760),
                      self.org, sz=13, c=self.C['text3'], b=True, align=PP_ALIGN.CENTER)

        # 클로징 푸터 (번호 없음)
        self._rect(s, 0, self.Y_FT, self.SW, self.H_FT, self.C['footer'])
        ft = ' | '.join(filter(None, ['Confidential', self.org, self.date]))
        self._txt(s, self.MX, self.Y_FT, Emu(7680960), self.H_FT,
                  ft, sz=8.5, c=self.C['text3'], align=PP_ALIGN.CENTER)

    # ── 저장 ──

    def save(self, path):
        """PPTX 파일 저장."""
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
